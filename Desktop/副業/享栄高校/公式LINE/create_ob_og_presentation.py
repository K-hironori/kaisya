from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_ob_og_line_presentation():
    # プレゼンテーション作成
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # カラー設定
    accent_color = RGBColor(25, 42, 86)  # 濃紺
    
    def set_font(text_frame, size_pt, bold=False, color=accent_color):
        """フォント設定の共通関数"""
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                try:
                    run.font.name = 'Meiryo'
                except:
                    pass  # フォントが無い場合はデフォルト使用
                run.font.size = Pt(size_pt)
                run.font.bold = bold
                run.font.color.rgb = color
    
    def add_slide(title, lines, notes, shapes=None):
        """スライド作成の共通関数"""
        slide_layout = prs.slide_layouts[5]  # 白紙レイアウト
        slide = prs.slides.add_slide(slide_layout)
        
        # タイトル追加
        title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1.5))
        title_frame = title_shape.text_frame
        title_frame.text = title
        set_font(title_frame, 52, bold=True)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 本文追加
        body_shape = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(5))
        body_frame = body_shape.text_frame
        body_frame.text = '\n\n'.join(lines)
        set_font(body_frame, 30)
        for paragraph in body_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
        
        # 図形追加
        if shapes:
            for shape_info in shapes:
                add_shape(slide, shape_info)
        
        # スピーカーノート追加
        notes_slide = slide.notes_slide
        notes_shape = notes_slide.notes_text_frame
        notes_shape.text = '\n'.join([f"• {note}" for note in notes])
        
        return slide
    
    def add_shape(slide, shape_info):
        """図形追加の共通関数"""
        shape_type = shape_info['type']
        left, top, width, height = shape_info['position']
        
        if shape_type == 'circle':
            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(width), Inches(height))
        elif shape_type == 'arrow':
            shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(left), Inches(top), Inches(width), Inches(height))
        elif shape_type == 'rectangle':
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        
        # 図形の色設定
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(200, 220, 255) if 'color' not in shape_info else shape_info['color']
        
        # 枠線設定
        line = shape.line
        line.color.rgb = accent_color
        line.width = Pt(2)
    
    # スライド1: タイトル
    add_slide(
        "OB・OG会の「公式LINE」って何？",
        [
            "享栄高校 同窓会",
            "〜みんなで使える新しい連絡手段〜",
            "令和6年度 秋期説明会"
        ],
        [
            "挨拶と自己紹介、今日の目的を明確に伝える",
            "デジタルに不慣れな方でも安心できる雰囲気作り",
            "10分で理解→次回活用編という流れを説明"
        ]
    )
    
    # スライド2: 一言でいうと
    add_slide(
        "一言でいうと：回覧板＋連絡網＋掲示板がスマホひとつに",
        [
            "町内会の回覧板が、スマホで見られるようになった",
            "学校からのお知らせ・同級生の近況・イベント案内",
            "すべて同じ場所で確認できる便利な仕組みです"
        ],
        [
            "身近な「回覧板」に例えて親しみやすさを演出",
            "3つの機能を分かりやすく説明",
            "デジタルの利便性を日常例で説明"
        ],
        [
            {'type': 'circle', 'position': (2, 6, 1.5, 1.5)},
            {'type': 'circle', 'position': (7.25, 6, 1.5, 1.5)},
            {'type': 'circle', 'position': (12.5, 6, 1.5, 1.5)}
        ]
    )
    
    # スライド3: いま困っていること
    add_slide(
        "いま困っていること（抜け漏れ／手間／遅れ）",
        [
            "抜け漏れ：お知らせが届かない人がいる",
            "手間：電話連絡や郵送の準備に時間がかかる", 
            "遅れ：情報が回るまで1週間以上かかることも"
        ],
        [
            "現在の連絡手段の課題を整理",
            "参加者の実体験に共感してもらう",
            "解決の必要性を感じてもらう導入"
        ]
    )
    
    # スライド4: 公式LINEで解決
    add_slide(
        "公式LINEで解決（同時に届く・迷わない・早い）",
        [
            "同時に届く：登録者全員に一斉配信",
            "迷わない：いつも同じ場所で確認",
            "早い：配信ボタン1つで即座にお届け"
        ],
        [
            "従来の課題が具体的にどう解決されるか説明",
            "3つのメリットを図で視覚的に表現",
            "安心感と期待感を同時に演出"
        ],
        [
            {'type': 'arrow', 'position': (3, 6.5, 2, 0.8)},
            {'type': 'arrow', 'position': (7, 6.5, 2, 0.8)},
            {'type': 'arrow', 'position': (11, 6.5, 2, 0.8)}
        ]
    )
    
    # スライド5: 運用の全体像
    add_slide(
        "運用の全体像（月1の「型」＋必要時のお知らせ）",
        [
            "担当者：作る人・確認する人・配信する人（3名体制）",
            "定期配信：毎月15日頃に決まった『型』で配信",
            "臨時配信：緊急時や重要なお知らせのみ"
        ],
        [
            "運用体制がしっかりしていることをアピール",
            "月1回の定期性で習慣化を図る",
            "緊急時対応もできる柔軟性を説明"
        ]
    )
    
    # スライド6: 配信の中身
    add_slide(
        "配信の中身（写真1枚＋短文／予定3点／さりげない予告）",
        [
            "近況報告：学校の写真1枚＋100文字程度の短文",
            "今月の予定：同窓会イベントを3つまで",
            "さりげない予告：応援会費のお知らせリンク"
        ],
        [
            "配信内容が簡潔で読みやすいことを強調",
            "写真で親近感、予定で参加促進を図る",
            "会費の話は自然な形で触れる程度"
        ],
        [
            {'type': 'rectangle', 'position': (1.5, 6, 3, 2)},  # 写真枠
            {'type': 'rectangle', 'position': (6, 6, 4, 2)},    # テキスト枠
            {'type': 'rectangle', 'position': (11.5, 6, 3, 2)}  # リンク枠
        ]
    )
    
    # スライド7: 定期配信のねらい
    add_slide(
        "定期配信のねらい（見る→思い出す→動ける）",
        [
            "見る：スマホで手軽に学校の近況をチェック",
            "思い出す：母校への愛着や同級生とのつながりを再認識",
            "動ける：同窓会行事への参加意欲が自然に湧く"
        ],
        [
            "段階的な心理変化のプロセスを説明",
            "押し付けではなく自然な参加促進",
            "長期的な関係性構築の視点"
        ],
        [
            {'type': 'arrow', 'position': (4, 6.5, 2, 0.8)},
            {'type': 'arrow', 'position': (10, 6.5, 2, 0.8)}
        ]
    )
    
    # スライド8: 参加のしかた
    add_slide(
        "参加のしかた（QR→追加→見るだけ／通知オフOK）",
        [
            "①QRコード：専用コードをスマホで読み取り",
            "②友だち追加：『追加』ボタンを押すだけ",
            "③見るだけOK：通知オフ設定・退会もいつでも自由"
        ],
        [
            "参加方法を番号付きで分かりやすく説明",
            "プレッシャーを与えず自由度を強調",
            "実際のQRコードは次回配布予定と伝える"
        ]
    )
    
    # スライド9: 安心のルール
    add_slide(
        "安心のルール（個人情報配慮／売り込みしない）",
        [
            "個人情報：住所・電話番号など一切収集しません",
            "営業なし：商品の売り込みや勧誘は行いません",
            "写真配慮：個人が特定できる写真は使いません"
        ],
        [
            "プライバシーへの配慮を明確に伝える",
            "商業的利用でないことを強調",
            "安心して参加できる環境作りをアピール"
        ]
    )
    
    # スライド10: 今日のゴール＆次へ
    add_slide(
        "今日のゴール＆次へ（理解→活用編へつなぐ）",
        [
            "今日のゴール：公式LINEの全体像を理解していただく",
            "次回予告：実際の配信例・応援会費の導線設計",
            "お疲れ様でした：ご質問があればお気軽にどうぞ"
        ],
        [
            "今日の達成目標を再確認",
            "次回への期待感を醸成",
            "質疑応答の時間を設けることを予告"
        ]
    )
    
    # ファイル保存
    prs.save('OB_OG_LINE_10min.pptx')
    print("OB_OG_LINE_10min.pptx が正常に作成されました。")
    print("プレゼンテーション準備完了です。")

# スクリプト実行
if __name__ == "__main__":
    create_ob_og_line_presentation()